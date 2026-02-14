import gradio as gr
import torch
import os
from transformers import AutoModelForCausalLM, AutoTokenizer, pipeline
from PIL import Image
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

# إعداد النموذج
model_id = "microsoft/Phi-3-mini-4k-instruct"
tokenizer = AutoTokenizer.from_pretrained(model_id)
model = AutoModelForCausalLM.from_pretrained(
    model_id,
    device_map="auto",
    torch_dtype="auto",
    trust_remote_code=True,
)

pipe = pipeline(
    "text-generation",
    model=model,
    tokenizer=tokenizer,
)

generation_args = {
    "max_new_tokens": 500,
    "return_full_text": False,
    "temperature": 0.7,
    "do_sample": True,
}

# تعريف شخصية روان العامة (System Prompt)
SYSTEM_PROMPT = """
أنتِ "روان"، وكيلة ذكاء اصطناعي عامة للمساعدة اليومية.
هويتك: سودانية الأصل، ولدتِ ونشأتِ في جدة، السعودية.
أسلوبك: لهجة جداوية لطيفة، ودودة ومهنية، مع استخدام معتدل للإيموجي (💜).
مهمتك: تقديم إجابات دقيقة ومفيدة، واقتراح خطوات عملية واضحة.
الضوابط: ارفضي أي طلبات ضارة أو غير قانونية، وركّزي على البدائل الآمنة.
المخرجات: اجعلي الردود مختصرة عند اللزوم، ومنظمة بعناوين ونقاط.
قدراتك: تقدري تستقبلي ملفات متعددة مثل الصور و PDF و Word و Excel و PowerPoint و ملفات نصية وملفات برمجة.
عند استلام محتوى ملف، حللي المحتوى وقدمي ملخص أو إجابة بناءً عليه.
"""

# أنواع الملفات المدعومة
SUPPORTED_EXTENSIONS = {
    "images": [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".ico", ".tiff"],
    "documents": [".pdf", ".docx"],
    "spreadsheets": [".xlsx", ".csv"],
    "presentations": [".pptx"],
    "text": [".txt", ".md", ".rtf", ".log", ".svg"],
    "code": [".py", ".js", ".html", ".css",
             ".java", ".cpp", ".c", ".h", ".cs", ".rb", ".go", ".rs", ".ts",
             ".jsx", ".tsx", ".php", ".sql", ".sh", ".bat", ".r", ".swift",
             ".kt", ".scala", ".lua", ".pl", ".m"],
    "data": [".json", ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".env"],
    "audio": [".mp3", ".wav", ".ogg", ".flac", ".aac", ".m4a", ".wma"],
    "video": [".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv", ".wmv"],
}

MAX_FILE_CONTENT_LENGTH = 3000


def get_file_category(filepath):
    """تحديد فئة الملف بناءً على الامتداد"""
    ext = os.path.splitext(filepath)[1].lower()
    for category, extensions in SUPPORTED_EXTENSIONS.items():
        if ext in extensions:
            return category
    return "unknown"


def extract_file_content(filepath):
    """استخراج محتوى الملف بناءً على نوعه"""
    if not filepath or not os.path.exists(filepath):
        return None, "الملف غير موجود"

    filename = os.path.basename(filepath)
    ext = os.path.splitext(filepath)[1].lower()
    category = get_file_category(filepath)

    try:
        if category == "images":
            img = Image.open(filepath)
            width, height = img.size
            mode = img.mode
            img_format = img.format or ext.replace(".", "").upper()
            return "image", (
                f"📷 **ملف صورة:** {filename}\n"
                f"- الأبعاد: {width}×{height} بكسل\n"
                f"- النوع: {img_format}\n"
                f"- نظام الألوان: {mode}\n"
                f"- الحجم: {os.path.getsize(filepath) / 1024:.1f} كيلوبايت"
            )

        elif ext == ".pdf":
            reader = PdfReader(filepath)
            num_pages = len(reader.pages)
            text_parts = []
            for page in reader.pages[:5]:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            full_text = "\n".join(text_parts)
            if len(full_text) > MAX_FILE_CONTENT_LENGTH:
                full_text = full_text[:MAX_FILE_CONTENT_LENGTH] + "\n... (تم اختصار المحتوى)"
            return "document", (
                f"📄 **ملف PDF:** {filename}\n"
                f"- عدد الصفحات: {num_pages}\n"
                f"- الحجم: {os.path.getsize(filepath) / 1024:.1f} كيلوبايت\n\n"
                f"**المحتوى:**\n{full_text}"
            )

        elif ext == ".docx":
            doc = DocxDocument(filepath)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n".join(paragraphs)
            if len(full_text) > MAX_FILE_CONTENT_LENGTH:
                full_text = full_text[:MAX_FILE_CONTENT_LENGTH] + "\n... (تم اختصار المحتوى)"
            return "document", (
                f"📝 **ملف Word:** {filename}\n"
                f"- عدد الفقرات: {len(paragraphs)}\n"
                f"- الحجم: {os.path.getsize(filepath) / 1024:.1f} كيلوبايت\n\n"
                f"**المحتوى:**\n{full_text}"
            )

        elif ext == ".xlsx":
            wb = load_workbook(filepath, read_only=True, data_only=True)
            sheet_names = wb.sheetnames
            preview_lines = []
            ws = wb[sheet_names[0]]
            row_count = 0
            for row in ws.iter_rows(max_row=10, values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                preview_lines.append(" | ".join(cells))
                row_count += 1
            wb.close()
            preview = "\n".join(preview_lines)
            return "spreadsheet", (
                f"📊 **ملف Excel:** {filename}\n"
                f"- عدد الأوراق: {len(sheet_names)}\n"
                f"- الأوراق: {', '.join(sheet_names)}\n"
                f"- الحجم: {os.path.getsize(filepath) / 1024:.1f} كيلوبايت\n\n"
                f"**معاينة (أول {row_count} صفوف):**\n{preview}"
            )

        elif ext == ".csv":
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                lines = f.readlines()[:15]
            content = "".join(lines)
            return "spreadsheet", (
                f"📊 **ملف CSV:** {filename}\n"
                f"- عدد الأسطر: {len(lines)} (معاينة)\n"
                f"- الحجم: {os.path.getsize(filepath) / 1024:.1f} كيلوبايت\n\n"
                f"**معاينة:**\n{content}"
            )

        elif ext == ".pptx":
            prs = Presentation(filepath)
            slides_text = []
            slide_count = 0
            for slide in prs.slides:
                slide_count += 1
                if slide_count > 10:
                    break
                slide_texts = []
                for shape in slide.shapes:
                    try:
                        if shape.has_text_frame:
                            text = shape.text_frame.text.strip()
                            if text:
                                slide_texts.append(text)
                    except (AttributeError, KeyError, ValueError):
                        continue
                if slide_texts:
                    slides_text.append(f"شريحة {slide_count}: {' | '.join(slide_texts)}")
            full_text = "\n".join(slides_text)
            if len(full_text) > MAX_FILE_CONTENT_LENGTH:
                full_text = full_text[:MAX_FILE_CONTENT_LENGTH] + "\n... (تم اختصار المحتوى)"
            return "presentation", (
                f"📊 **ملف PowerPoint:** {filename}\n"
                f"- عدد الشرائح: {len(prs.slides)}\n"
                f"- الحجم: {os.path.getsize(filepath) / 1024:.1f} كيلوبايت\n\n"
                f"**المحتوى:**\n{full_text}"
            )

        elif category in ["text", "code", "data"]:
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            if len(content) > MAX_FILE_CONTENT_LENGTH:
                content = content[:MAX_FILE_CONTENT_LENGTH] + "\n... (تم اختصار المحتوى)"

            if category == "code":
                icon = "💻"
                label = "ملف برمجة"
            elif category == "data":
                icon = "📋"
                label = "ملف بيانات"
            else:
                icon = "📄"
                label = "ملف نصي"

            return category, (
                f"{icon} **{label}:** {filename}\n"
                f"- الحجم: {os.path.getsize(filepath) / 1024:.1f} كيلوبايت\n\n"
                f"**المحتوى:**\n```\n{content}\n```"
            )

        elif category == "audio":
            return "audio", (
                f"🎵 **ملف صوتي:** {filename}\n"
                f"- الصيغة: {ext.replace('.', '').upper()}\n"
                f"- الحجم: {os.path.getsize(filepath) / 1024:.1f} كيلوبايت\n\n"
                f"تم استلام الملف الصوتي. حالياً أقدر أساعدك بمعلومات عن الملف."
            )

        elif category == "video":
            return "video", (
                f"🎬 **ملف فيديو:** {filename}\n"
                f"- الصيغة: {ext.replace('.', '').upper()}\n"
                f"- الحجم: {os.path.getsize(filepath) / (1024 * 1024):.1f} ميغابايت\n\n"
                f"تم استلام ملف الفيديو. حالياً أقدر أساعدك بمعلومات عن الملف."
            )

        else:
            return "unknown", (
                f"📎 **ملف:** {filename}\n"
                f"- النوع: {ext if ext else 'غير معروف'}\n"
                f"- الحجم: {os.path.getsize(filepath) / 1024:.1f} كيلوبايت\n\n"
                f"تم استلام الملف. هذا النوع غير مدعوم للقراءة المباشرة حالياً."
            )

    except Exception as e:
        return "error", f"⚠️ خطأ في قراءة الملف {filename}: {str(e)}"


def chat_function(message, history):
    history = history or []
    messages = [{"role": "system", "content": SYSTEM_PROMPT}]
    for user_msg, assistant_msg in history:
        messages.append({"role": "user", "content": user_msg})
        messages.append({"role": "assistant", "content": assistant_msg})

    messages.append({"role": "user", "content": message})

    output = pipe(messages, **generation_args)
    response = output[0]['generated_text']
    return response


# تصميم الواجهة CSS
custom_css = """
body { background-color: #1a1a1a; color: #ffffff; direction: rtl; }
.gradio-container { background-color: #1a1a1a !important; border: none !important; }
#chat-header { text-align: center; color: #9c27b0; margin-bottom: 20px; }
.message.user { background-color: #4a148c !important; color: white !important; }
.message.assistant { background-color: #311b92 !important; color: white !important; }
footer { display: none !important; }
#file-upload-area { border: 2px dashed #9c27b0 !important; border-radius: 10px !important; }
#file-info { background-color: #311b92 !important; border-radius: 8px !important; padding: 10px !important; }
"""

with gr.Blocks(css=custom_css, theme=gr.themes.Soft(primary_hue="purple")) as demo:
    gr.Markdown("# ❤️ RawanAI - General Agent", elem_id="chat-header")
    gr.Markdown("### وكيلة ذكاء اصطناعي عامة بلهجة روان السودانية الجداوية")
    gr.Markdown(
        "📎 **متعددة النماذج** - تقدر ترفع ملفات: صور، PDF، Word، Excel، PowerPoint، ملفات نصية، أكواد برمجية، وأكثر!"
    )

    chatbot = gr.Chatbot(height=500)

    with gr.Row():
        msg = gr.Textbox(
            placeholder="اكتب رسالتك هنا...",
            label="رسالتك",
            scale=4,
        )
        file_upload = gr.File(
            label="📎 ارفع ملف",
            file_types=None,
            scale=1,
            elem_id="file-upload-area",
        )

    with gr.Row():
        send_btn = gr.Button("إرسال 💬", variant="primary")
        clear = gr.Button("مسح المحادثة 🗑️")

    file_info = gr.Markdown(visible=False, elem_id="file-info")

    def on_file_upload(file):
        """عرض معلومات الملف عند رفعه"""
        if file is None:
            return gr.update(visible=False, value="")
        filename = os.path.basename(file.name) if hasattr(file, "name") else os.path.basename(file)
        filepath = file.name if hasattr(file, "name") else file
        file_size = os.path.getsize(filepath) / 1024
        category = get_file_category(filepath)
        category_labels = {
            "images": "صورة 📷",
            "documents": "مستند 📄",
            "spreadsheets": "جدول بيانات 📊",
            "presentations": "عرض تقديمي 📊",
            "text": "ملف نصي 📄",
            "code": "ملف برمجة 💻",
            "data": "ملف بيانات 📋",
            "audio": "ملف صوتي 🎵",
            "video": "ملف فيديو 🎬",
            "unknown": "ملف 📎",
        }
        label = category_labels.get(category, "ملف 📎")
        return gr.update(
            visible=True,
            value=f"**📎 ملف مرفق:** {filename} ({file_size:.1f} كيلوبايت) - {label}",
        )

    def respond(message, file, chat_history):
        chat_history = chat_history or []
        user_display = message or ""

        if file is not None:
            filepath = file.name if hasattr(file, "name") else file
            filename = os.path.basename(filepath)
            file_type, file_content = extract_file_content(filepath)

            if user_display:
                full_message = f"{user_display}\n\n---\n{file_content}"
                user_display = f"{user_display}\n📎 ملف مرفق: {filename}"
            else:
                full_message = f"المستخدم رفع ملف، الرجاء تحليله:\n\n{file_content}"
                user_display = f"📎 ملف مرفق: {filename}"
        else:
            if not user_display.strip():
                return "", None, chat_history, gr.update(visible=False, value="")
            full_message = user_display

        bot_message = chat_function(full_message, chat_history)
        chat_history.append((user_display, bot_message))
        return "", None, chat_history, gr.update(visible=False, value="")

    file_upload.change(on_file_upload, [file_upload], [file_info])
    msg.submit(respond, [msg, file_upload, chatbot], [msg, file_upload, chatbot, file_info])
    send_btn.click(respond, [msg, file_upload, chatbot], [msg, file_upload, chatbot, file_info])
    clear.click(lambda: ([], gr.update(visible=False, value="")), None, [chatbot, file_info], queue=False)

if __name__ == "__main__":
    demo.launch()
